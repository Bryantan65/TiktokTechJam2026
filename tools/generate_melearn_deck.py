"""Generate the editable MeLearn three-minute demo deck.

All visuals are native PowerPoint shapes so the resulting deck can be adjusted
without access to the generator. Speaker notes contain the timed narration.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs" / "MeLearn-Demo-Deck.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = "071426"
NAVY_2 = "0C2038"
PANEL = "102A46"
PANEL_2 = "163754"
BLUE = "3B82F6"
BLUE_DARK = "1D4ED8"
YELLOW = "FBBF24"
YELLOW_DARK = "D99A00"
WHITE = "F8FAFC"
MUTED = "A9B8C9"
FAINT = "6F86A0"
GREEN = "34D399"
ROSE = "FB7185"

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"
FONT_MONO = "Consolas"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def set_background(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_rect(slide, x, y, w, h, fill, radius=True, line=None, line_width=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=WHITE,
    bold=False,
    font=FONT_BODY,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    # Newlines can be represented as multiple runs by presentation renderers.
    # Style every run so wrapped/manual-break lines retain the intended color.
    for run in paragraph.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_rich_line(slide, runs, x, y, w, h, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for spec in runs:
        run = paragraph.add_run()
        run.text = spec["text"]
        run.font.name = spec.get("font", FONT_BODY)
        run.font.size = Pt(spec.get("size", 18))
        run.font.bold = spec.get("bold", False)
        run.font.color.rgb = rgb(spec.get("color", WHITE))
    return box


def add_brand(slide, x, y, w=3.0, h=0.5, size=23, align=PP_ALIGN.LEFT):
    return add_rich_line(
        slide,
        [
            {"text": "M", "color": BLUE, "bold": True, "size": size, "font": FONT_HEAD},
            {"text": "e", "color": YELLOW, "bold": True, "size": size, "font": FONT_HEAD},
            {"text": "L", "color": BLUE, "bold": True, "size": size, "font": FONT_HEAD},
            {"text": "earn", "color": YELLOW, "bold": True, "size": size, "font": FONT_HEAD},
        ],
        x,
        y,
        w,
        h,
        align=align,
    )


def add_kicker(slide, text, x=0.68, y=0.42, w=6.5):
    add_text(slide, text.upper(), x, y, w, 0.24, 9, YELLOW, True, FONT_BODY)


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.68, 0.48, 12.0, 0.58, 27, WHITE, True, FONT_HEAD)
    if subtitle:
        add_text(slide, subtitle, 0.7, 1.10, 11.9, 0.42, 12, MUTED)


def add_footer(slide, number):
    add_brand(slide, 0.68, 7.14, 1.6, 0.22, 10)
    add_text(slide, f"{number:02d}", 12.16, 7.13, 0.48, 0.2, 9, FAINT, True, align=PP_ALIGN.RIGHT)
    add_rect(slide, 2.25, 7.22, 9.55, 0.018, PANEL_2, radius=False)
    add_rect(slide, 2.25, 7.22, 9.55 * number / 6, 0.018, BLUE, radius=False)


def add_notes(slide, text):
    frame = slide.notes_slide.notes_text_frame
    frame.text = text.strip()


def add_pill(slide, text, x, y, w, fill=PANEL_2, color=MUTED, line=None, size=10):
    add_rect(slide, x, y, w, 0.34, fill, True, line)
    add_text(slide, text, x, y + 0.01, w, 0.28, size, color, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_metric_card(slide, x, y, w, label, value, accent=YELLOW, note=None):
    add_rect(slide, x, y, w, 1.02, PANEL, True, PANEL_2)
    add_text(slide, label.upper(), x + 0.18, y + 0.16, w - 0.36, 0.2, 8.5, MUTED, True)
    add_text(slide, value, x + 0.18, y + 0.40, w - 0.36, 0.36, 20, accent, True, FONT_HEAD)
    if note:
        add_text(slide, note, x + 0.18, y + 0.77, w - 0.36, 0.17, 8, FAINT)


def add_chevron(slide, x, y, color=FAINT, size=20):
    add_text(slide, ">", x, y, 0.3, 0.35, size, color, True, align=PP_ALIGN.CENTER)


def slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    # Subtle edge accents.
    add_rect(slide, 0, 0, 0.09, SLIDE_H, BLUE, radius=False)
    add_rect(slide, 0.09, 0, 0.035, 2.7, YELLOW, radius=False)
    add_text(slide, "AUTONOMOUS RECOMMENDER RESEARCH", 0.7, 0.55, 5.8, 0.3, 10, YELLOW, True)
    add_brand(slide, 0.68, 1.02, 6.5, 0.86, 48)
    add_text(
        slide,
        "Machine learning that learns\nhow to improve itself.",
        0.7,
        2.02,
        6.6,
        1.28,
        27,
        WHITE,
        True,
        FONT_HEAD,
        line_spacing=0.9,
    )
    add_text(
        slide,
        "An autonomous agent that writes, evaluates, and improves\nrecommendation models—without a human in the loop.",
        0.72,
        3.48,
        6.35,
        0.75,
        13,
        MUTED,
        line_spacing=1.05,
    )

    # Hero result card.
    add_rect(slide, 8.0, 0.72, 4.62, 4.22, PANEL, True, PANEL_2, 1.2)
    add_pill(slide, "KUAIRAND · REQUIRED", 8.35, 1.04, 1.92, PANEL_2, MUTED)
    add_text(slide, "+0.003944", 8.35, 1.70, 3.95, 0.76, 35, YELLOW, True, FONT_HEAD)
    add_text(slide, "PRIMARY SCORE GAIN", 8.37, 2.44, 3.45, 0.25, 9, MUTED, True)
    add_rect(slide, 8.37, 3.04, 3.72, 0.13, PANEL_2, True)
    add_rect(slide, 8.37, 3.04, 3.26, 0.13, BLUE, True)
    add_rect(slide, 10.02, 2.87, 0.035, 0.47, YELLOW, radius=False)
    add_text(slide, "target +0.002", 9.46, 3.42, 1.18, 0.22, 8, YELLOW, True, align=PP_ALIGN.CENTER)
    add_text(slide, "MeLearn cleared the required margin\non validation—using the official evaluator.", 8.37, 3.92, 3.72, 0.58, 11, WHITE)

    add_metric_card(slide, 0.7, 5.35, 2.62, "experiments", "30", BLUE, "of a 50-experiment cap")
    add_metric_card(slide, 3.52, 5.35, 2.62, "total cost", "$3.65", YELLOW, "2.46M model tokens")
    add_metric_card(slide, 6.34, 5.35, 2.62, "human interventions", "0", GREEN, "after the run started")
    add_metric_card(slide, 9.16, 5.35, 3.46, "GPU-hours", "0", BLUE, "CPU-only record run")
    add_text(slide, "TIKTOK TECHJAM 2026 · DEMO", 0.7, 6.76, 4.4, 0.22, 8.5, FAINT, True)

    add_notes(
        slide,
        "MeLearn is machine learning that learns how to improve itself. With one command, it ran thirty experiments, beat the required KuaiRand baseline, cost three dollars and sixty-five cents, and required zero human interventions.",
    )
    return slide


def slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "Improving a recommender is still a manual loop.", "Tiny ranking gains demand disciplined experimentation—not one lucky training run.")

    stages = [
        ("01", "Hypothesize", "Choose what\nto change"),
        ("02", "Code", "Build a full\nsolution"),
        ("03", "Train", "Run the\nexperiment"),
        ("04", "Evaluate", "Measure three\nrandom seeds"),
        ("05", "Decide", "Keep, reject,\nor branch"),
    ]
    x0, gap, card_w = 0.7, 0.24, 2.25
    for i, (num, label, desc) in enumerate(stages):
        x = x0 + i * (card_w + gap)
        add_rect(slide, x, 1.88, card_w, 1.52, PANEL, True, PANEL_2)
        add_text(slide, num, x + 0.16, 2.06, 0.34, 0.24, 9, BLUE, True)
        add_text(slide, label, x + 0.16, 2.35, card_w - 0.32, 0.33, 15, WHITE, True, FONT_HEAD)
        add_text(slide, desc, x + 0.16, 2.76, card_w - 0.32, 0.46, 10, MUTED)
        if i < len(stages) - 1:
            add_chevron(slide, x + card_w + 0.01, 2.48, BLUE, 18)

    add_text(slide, "The benchmark makes that loop unforgiving", 0.7, 3.86, 6.4, 0.35, 17, WHITE, True, FONT_HEAD)
    challenge_cards = [
        ("TASK", "Rank long_view", "within each user's impressions", BLUE),
        ("TARGET", "+0.002", "above the official FM baseline", YELLOW),
        ("BUDGET", "50 exp · 6 h", "hard competition cap", BLUE),
        ("NOISE", "3 seeds", "one seed can swing 65% of target", GREEN),
    ]
    for i, (label, value, note, accent) in enumerate(challenge_cards):
        x = 0.7 + i * 3.12
        add_rect(slide, x, 4.42, 2.88, 1.30, PANEL, True, PANEL_2)
        add_text(slide, label, x + 0.18, 4.60, 2.5, 0.18, 8.5, accent, True)
        add_text(slide, value, x + 0.18, 4.88, 2.5, 0.30, 16, WHITE, True, FONT_HEAD)
        add_text(slide, note, x + 0.18, 5.29, 2.50, 0.24, 9, MUTED)

    add_rect(slide, 0.7, 6.08, 11.92, 0.72, PANEL_2, True)
    add_text(slide, "MeLearn automates the whole research loop—not just model inference.", 0.98, 6.26, 11.36, 0.30, 15, YELLOW, True, FONT_HEAD, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)
    add_notes(
        slide,
        "Recommendation research is an iterative loop: form a hypothesis, write the code, train, evaluate, and decide what comes next. Here, the agent had to rank watch completion, beat the official factorization-machine baseline by point zero zero two, and stay inside fifty experiments and six hours. Because seed noise can consume most of that margin, a single promising score is not enough.",
    )
    return slide


def terminal_line(slide, text, x, y, color=MUTED, bold=False, size=9.5):
    add_text(slide, text, x, y, 6.6, 0.25, size, color, bold, FONT_MONO)


def slide_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "One command. Then leave it.", "A multi-hour autonomous run, compressed here into one visible experiment cycle.")

    # Terminal mockup.
    add_rect(slide, 0.68, 1.68, 7.35, 4.90, "08111F", True, PANEL_2)
    add_rect(slide, 0.68, 1.68, 7.35, 0.43, PANEL_2, True)
    for i, c in enumerate([ROSE, YELLOW, GREEN]):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.92 + i * 0.24), Inches(1.82), Inches(0.10), Inches(0.10))
        circle.fill.solid(); circle.fill.fore_color.rgb = rgb(c); circle.line.fill.background()
    add_text(slide, "melearn · autonomous run", 1.85, 1.78, 4.9, 0.20, 8.5, MUTED, True, FONT_MONO, align=PP_ALIGN.CENTER)

    lines = [
        ("> python -m agent --dataset pure", BLUE, True),
        ("[run] ledger loaded · 26 completed experiments", MUTED, False),
        ("[uct] selected parent node 24", YELLOW, False),
        ("[idea] add a small DeepFM ensemble member", WHITE, False),
        ("[write] solutions/027_deepfm_ensemble.py", MUTED, False),
        ("[eval] seed 1/3  seed 2/3  seed 3/3", BLUE, False),
        ("[score] primary 0.605493  ·  gain +0.003944", GREEN, True),
        ("[ledger] verdict: KEPT  ·  continue search", YELLOW, True),
    ]
    for i, (text, color, bold) in enumerate(lines):
        terminal_line(slide, text, 1.02, 2.35 + i * 0.45, color, bold)
    add_pill(slide, "HOURS COMPRESSED FOR DEMO", 1.02, 6.00, 2.36, PANEL, FAINT, size=8)

    # Right-hand experiment cycle.
    add_rect(slide, 8.34, 1.68, 4.29, 4.90, PANEL, True, PANEL_2)
    add_text(slide, "WHAT THE AGENT DOES", 8.65, 1.98, 3.64, 0.22, 9, YELLOW, True)
    steps = [
        ("01", "Read the research history"),
        ("02", "Choose a promising branch"),
        ("03", "Write a standalone solution"),
        ("04", "Evaluate across three seeds"),
        ("05", "Persist the result and repeat"),
    ]
    for i, (num, label) in enumerate(steps):
        y = 2.40 + i * 0.61
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.66), Inches(y), Inches(0.34), Inches(0.34))
        circle.fill.solid(); circle.fill.fore_color.rgb = rgb(BLUE if i < 4 else YELLOW); circle.line.fill.background()
        add_text(slide, num, 8.66, y + 0.06, 0.34, 0.14, 7.5, NAVY, True, align=PP_ALIGN.CENTER)
        add_text(slide, label, 9.18, y + 0.01, 3.05, 0.28, 10.5, WHITE, i == 4)
        if i < 4:
            add_rect(slide, 8.81, y + 0.36, 0.025, 0.24, PANEL_2, radius=False)
    add_rect(slide, 8.64, 5.72, 3.66, 0.56, "183D38", True, GREEN)
    add_text(slide, "WINNER · EXPERIMENT 27", 8.82, 5.86, 3.30, 0.23, 10, GREEN, True, align=PP_ALIGN.CENTER)

    add_footer(slide, 3)
    add_notes(
        slide,
        "We start MeLearn with one command. It reads every previous experiment from the ledger, then uses tree search to select a promising parent. It states a hypothesis, writes a complete Python solution, and hands that code to the harness. The harness evaluates three seeds and returns only the validation result. MeLearn records the score, keeps or rejects the idea, and starts again. This footage compresses hours into seconds; after twenty-seven experiments, the highlighted branch produced the winning model.",
    )
    return slide


def add_loop_card(slide, x, y, num, title, detail, accent=BLUE):
    add_rect(slide, x, y, 2.18, 1.02, PANEL, True, PANEL_2)
    add_text(slide, num, x + 0.15, y + 0.15, 0.30, 0.18, 8, accent, True)
    add_text(slide, title, x + 0.52, y + 0.13, 1.48, 0.23, 11, WHITE, True, FONT_HEAD)
    add_text(slide, detail, x + 0.15, y + 0.50, 1.85, 0.34, 8.7, MUTED)


def slide_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "The harness makes autonomy trustworthy.", "The model proposes experiments; the surrounding system owns truth, limits, and memory.")

    # Search loop on the left.
    add_text(slide, "THE CLOSED RESEARCH LOOP", 0.7, 1.68, 6.8, 0.22, 9, YELLOW, True)
    cards = [
        (0.70, 2.05, "01", "Ledger", "Persistent run memory", YELLOW),
        (3.04, 2.05, "02", "UCT select", "Explore or exploit", BLUE),
        (5.38, 2.05, "03", "Hypothesis", "State the next idea", BLUE),
        (5.38, 3.57, "04", "Write", "Complete Python solution", BLUE),
        (3.04, 3.57, "05", "Evaluate", "Official metric · 3 seeds", GREEN),
        (0.70, 3.57, "06", "Verdict", "Keep, reject, or branch", YELLOW),
    ]
    for card in cards:
        add_loop_card(slide, *card)
    add_chevron(slide, 2.75, 2.36, BLUE, 20)
    add_chevron(slide, 5.09, 2.36, BLUE, 20)
    add_text(slide, "v", 6.28, 3.19, 0.3, 0.25, 16, BLUE, True, align=PP_ALIGN.CENTER)
    add_text(slide, "<", 5.03, 3.88, 0.3, 0.25, 20, BLUE, True, align=PP_ALIGN.CENTER)
    add_text(slide, "<", 2.69, 3.88, 0.3, 0.25, 20, BLUE, True, align=PP_ALIGN.CENTER)
    add_text(slide, "^", 1.63, 3.19, 0.3, 0.25, 16, YELLOW, True, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.70, 5.05, 6.86, 0.83, PANEL_2, True)
    add_rich_line(
        slide,
        [
            {"text": "The ledger is the memory. ", "color": YELLOW, "bold": True, "size": 13, "font": FONT_HEAD},
            {"text": "The model starts every turn clean.", "color": WHITE, "bold": True, "size": 13, "font": FONT_HEAD},
        ],
        0.99,
        5.28,
        6.27,
        0.28,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, "UCT can return to an idea abandoned twenty iterations ago.", 1.02, 5.60, 6.18, 0.20, 9, MUTED, align=PP_ALIGN.CENTER)

    # Guardrails panel on the right.
    add_rect(slide, 8.02, 1.68, 4.60, 4.91, PANEL, True, PANEL_2)
    add_text(slide, "MECHANICAL GUARDRAILS", 8.35, 1.98, 3.94, 0.22, 9, YELLOW, True)
    guardrails = [
        ("TEST", "Labels are never exposed"),
        ("SCORE", "Official evaluator stays read-only"),
        ("SEEDS", "Every candidate runs three times"),
        ("INPUT", "Invalid or duplicate outputs are blocked"),
        ("FAIL", "Crashes return as actionable text"),
    ]
    for i, (tag, label) in enumerate(guardrails):
        y = 2.48 + i * 0.70
        add_pill(slide, tag, 8.35, y, 0.72, PANEL_2, BLUE if i < 4 else ROSE, size=7.5)
        add_text(slide, label, 9.28, y + 0.03, 2.96, 0.28, 10.3, WHITE, i in (0, 1))
    add_rect(slide, 8.35, 6.03, 3.94, 0.32, "183D38", True)
    add_text(slide, "The agent cannot reason around these rules.", 8.48, 6.09, 3.68, 0.16, 8, GREEN, True, align=PP_ALIGN.CENTER)

    add_footer(slide, 4)
    add_notes(
        slide,
        "Autonomy is only useful if the evaluation is trustworthy. The ledger is persistent memory, and UCT lets MeLearn revisit older branches instead of endlessly refining its latest idea. The harness—not the model—owns the data, official evaluator, seed policy, and time limits. Test labels are never visible, invalid or duplicate solutions are rejected, and failures return as text so the agent can diagnose and recover on its next turn.",
    )
    return slide


def add_result_panel(slide, x, y, w, dataset, badge, baseline, score, gain, max_gain, metrics, target=None):
    add_rect(slide, x, y, w, 3.80, PANEL, True, PANEL_2)
    add_text(slide, dataset, x + 0.28, y + 0.28, w - 2.0, 0.34, 17, WHITE, True, FONT_HEAD)
    add_pill(slide, badge, x + w - 1.52, y + 0.27, 1.20, PANEL_2, YELLOW if badge == "REQUIRED" else BLUE, size=8)
    add_text(slide, "BASELINE", x + 0.28, y + 0.86, 1.15, 0.18, 8, MUTED, True)
    add_text(slide, baseline, x + 0.28, y + 1.08, 1.60, 0.40, 19, WHITE, True, FONT_HEAD)
    add_text(slide, "MELEARN", x + 2.18, y + 0.86, 1.15, 0.18, 8, BLUE, True)
    add_text(slide, score, x + 2.18, y + 1.08, 2.20, 0.40, 19, BLUE, True, FONT_HEAD)
    add_text(slide, "GAIN", x + w - 1.55, y + 0.86, 1.25, 0.18, 8, YELLOW, True, align=PP_ALIGN.RIGHT)
    add_text(slide, gain, x + w - 2.15, y + 1.08, 1.85, 0.40, 19, YELLOW, True, FONT_HEAD, align=PP_ALIGN.RIGHT)

    # Gain track; scale is explicit within each panel.
    track_x, track_y, track_w = x + 0.30, y + 1.82, w - 0.60
    add_text(slide, "GAIN OVER BASELINE", track_x, track_y - 0.24, 2.1, 0.18, 7.5, FAINT, True)
    add_rect(slide, track_x, track_y, track_w, 0.14, PANEL_2, True)
    ratio = min(float(gain.replace("+", "")) / max_gain, 1.0)
    add_rect(slide, track_x, track_y, track_w * ratio, 0.14, YELLOW, True)
    if target is not None:
        target_x = track_x + track_w * target / max_gain
        add_rect(slide, target_x, track_y - 0.09, 0.025, 0.34, WHITE, radius=False)
        add_text(slide, "target", target_x - 0.28, track_y + 0.29, 0.58, 0.15, 6.5, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, "0", track_x, track_y + 0.19, 0.3, 0.16, 7, FAINT)
    add_text(slide, f"+{max_gain:.3f}", track_x + track_w - 0.65, track_y + 0.19, 0.65, 0.16, 7, FAINT, align=PP_ALIGN.RIGHT)

    # Metric-specific deltas.
    for i, (name, value) in enumerate(metrics):
        mx = x + 0.30 + i * ((w - 0.78) / 2)
        mw = (w - 0.90) / 2
        add_rect(slide, mx, y + 2.48, mw, 0.88, NAVY_2, True)
        add_text(slide, name, mx + 0.16, y + 2.63, mw - 0.32, 0.18, 8, MUTED, True)
        add_text(slide, value, mx + 0.16, y + 2.91, mw - 0.32, 0.25, 14, GREEN, True, FONT_HEAD)


def slide_five(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "MeLearn cleared the required target.", "Validation-best results · primary score = mean(GAUC, nDCG@5) · official evaluator unmodified")

    add_result_panel(
        slide,
        0.68,
        1.70,
        5.88,
        "KuaiRand-Pure",
        "REQUIRED",
        "0.6016",
        "0.605493",
        "+0.003944",
        0.0045,
        [("GAUC DELTA", "+0.005069"), ("nDCG@5 DELTA", "+0.002818")],
        0.002,
    )
    add_result_panel(
        slide,
        6.78,
        1.70,
        5.85,
        "KuaiRand-1K",
        "BONUS",
        "0.6451",
        "0.682994",
        "+0.037894",
        0.040,
        [("GAUC DELTA", "+0.027823"), ("nDCG@5 DELTA", "+0.047965")],
    )

    stats = [
        ("EXPERIMENTS", "30 / 50", BLUE),
        ("TOKENS", "2.46M", YELLOW),
        ("TOTAL COST", "$3.65", YELLOW),
        ("GPU-HOURS", "0", GREEN),
        ("HUMAN INPUT", "0", GREEN),
    ]
    for i, (label, value, accent) in enumerate(stats):
        x = 0.68 + i * 2.43
        add_rect(slide, x, 5.78, 2.21, 0.86, PANEL_2, True)
        add_text(slide, label, x + 0.13, 5.92, 1.95, 0.16, 7.5, MUTED, True, align=PP_ALIGN.CENTER)
        add_text(slide, value, x + 0.13, 6.19, 1.95, 0.25, 14, accent, True, FONT_HEAD, align=PP_ALIGN.CENTER)

    add_footer(slide, 5)
    add_notes(
        slide,
        "On the required KuaiRand-Pure dataset, the official baseline scored point six zero one six. MeLearn reached point six zero five four nine three: a gain of point zero zero three nine four four, almost twice the required margin. On the bonus one-thousand-user dataset, it improved the primary score by point zero three seven eight nine four. The record run used thirty of fifty experiments, two point four six million tokens, three dollars and sixty-five cents, no GPUs, and no human intervention.",
    )
    return slide


def add_comparison_bar(slide, x, y, w, label, value_text, fraction, accent):
    add_text(slide, label, x, y, 0.62, 0.20, 9, MUTED, True)
    add_rect(slide, x + 0.72, y + 0.03, w - 1.72, 0.18, PANEL_2, True)
    add_rect(slide, x + 0.72, y + 0.03, (w - 1.72) * fraction, 0.18, accent, True)
    add_text(slide, value_text, x + w - 0.86, y - 0.03, 0.86, 0.24, 11, accent, True, FONT_HEAD, align=PP_ALIGN.RIGHT)


def slide_six(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, "About 10× the overlap. About 10× the gain.", "The same agent and harness exposed where the recommendation signal can—and cannot—transfer.")

    add_rect(slide, 0.68, 1.72, 5.60, 2.66, PANEL, True, PANEL_2)
    add_text(slide, "SEEN USER–CREATOR PAIRS", 0.98, 2.02, 4.98, 0.22, 9, BLUE, True)
    add_text(slide, "Transferable interaction history", 0.98, 2.35, 4.98, 0.32, 16, WHITE, True, FONT_HEAD)
    add_comparison_bar(slide, 0.98, 3.02, 4.98, "Pure", "3.38%", 3.38 / 33.70, BLUE)
    add_comparison_bar(slide, 0.98, 3.56, 4.98, "1K", "33.70%", 1.0, YELLOW)

    add_rect(slide, 7.03, 1.72, 5.60, 2.66, PANEL, True, PANEL_2)
    add_text(slide, "PRIMARY SCORE DELTA", 7.33, 2.02, 4.98, 0.22, 9, YELLOW, True)
    add_text(slide, "Improvement over each baseline", 7.33, 2.35, 4.98, 0.32, 16, WHITE, True, FONT_HEAD)
    add_comparison_bar(slide, 7.33, 3.02, 4.98, "Pure", "+.0039", 0.003944 / 0.037894, BLUE)
    add_comparison_bar(slide, 7.33, 3.56, 4.98, "1K", "+.0379", 1.0, YELLOW)

    # The title carries the 10x comparison; keep the plot area unobstructed.
    add_text(slide, "≈", 6.39, 2.92, 0.54, 0.30, 18, MUTED, True, FONT_HEAD, align=PP_ALIGN.CENTER)
    add_text(slide, "10×", 6.31, 3.28, 0.70, 0.34, 15, YELLOW, True, FONT_HEAD, align=PP_ALIGN.CENTER)

    add_rect(slide, 0.68, 4.83, 11.95, 1.46, PANEL_2, True)
    add_text(slide, "MeLearn did not just find a better model.", 0.98, 5.15, 11.35, 0.36, 20, WHITE, True, FONT_HEAD, align=PP_ALIGN.CENTER)
    add_text(slide, "It found where the learnable signal lives.", 0.98, 5.61, 11.35, 0.38, 21, YELLOW, True, FONT_HEAD, align=PP_ALIGN.CENTER)

    add_brand(slide, 0.68, 6.70, 2.4, 0.35, 17)
    add_text(slide, "Measures · learns · recovers", 3.04, 6.75, 4.3, 0.20, 9, MUTED, True)
    add_text(slide, "AUTONOMOUS RECOMMENDER RESEARCH", 8.20, 6.73, 4.42, 0.20, 8.5, FAINT, True, align=PP_ALIGN.RIGHT)
    add_footer(slide, 6)
    add_notes(
        slide,
        "The largest result also revealed the central insight. In Pure, only three point three eight percent of validation rows contain a user-creator pair seen during training. In one-K, that overlap is thirty-three point seven percent—about ten times larger. The performance gain grows by almost the same factor. MeLearn did not just find a better model. It discovered where the dataset's learnable signal lives: an autonomous researcher that measures, learns, and recovers.",
    )
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "MeLearn — Three-Minute Demo Deck"
    prs.core_properties.subject = "TikTok TechJam 2026 autonomous recommender research agent"
    prs.core_properties.author = "MeLearn team"
    prs.core_properties.keywords = "MeLearn, autonomous agent, recommendation, KuaiRand, TechJam"
    prs.core_properties.comments = "Generated from editable native PowerPoint shapes."

    slide_one(prs)
    slide_two(prs)
    slide_three(prs)
    slide_four(prs)
    slide_five(prs)
    slide_six(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build_deck())
